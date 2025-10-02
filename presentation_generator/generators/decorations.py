import random
import sys
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR, MSO_FILL
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import math

def get_resource_path(relative_path):
    if getattr(sys, 'frozen', False):
        return Path(sys._MEIPASS) / relative_path
    else:
        return Path(__file__).parent.parent.parent / relative_path


class SlideDecorator:
    def __init__(self):
        from ..localization.manager import get_localization_manager
        self.loc = get_localization_manager()
        self.color_schemes = {
            'modern_blue': {
                'primary': RGBColor(41, 98, 255),
                'secondary': RGBColor(116, 185, 255),
                'accent': RGBColor(255, 107, 107),
                'background': RGBColor(248, 250, 252),
                'text': RGBColor(30, 41, 59)
            },
            'elegant_purple': {
                'primary': RGBColor(139, 69, 255),
                'secondary': RGBColor(196, 181, 253),
                'accent': RGBColor(255, 154, 0),
                'background': RGBColor(250, 249, 255),
                'text': RGBColor(55, 48, 163)
            },
            'warm_orange': {
                'primary': RGBColor(255, 107, 0),
                'secondary': RGBColor(255, 183, 77),
                'accent': RGBColor(34, 197, 94),
                'background': RGBColor(255, 251, 235),
                'text': RGBColor(124, 45, 18)
            },
            'forest_green': {
                'primary': RGBColor(34, 197, 94),
                'secondary': RGBColor(134, 239, 172),
                'accent': RGBColor(239, 68, 68),
                'background': RGBColor(240, 253, 244),
                'text': RGBColor(20, 83, 45)
            },
            'ocean_teal': {
                'primary': RGBColor(20, 184, 166),
                'secondary': RGBColor(153, 246, 228),
                'accent': RGBColor(251, 146, 60),
                'background': RGBColor(240, 253, 250),
                'text': RGBColor(19, 78, 74)
            },
            'sunset_pink': {
                'primary': RGBColor(236, 72, 153),
                'secondary': RGBColor(251, 207, 232),
                'accent': RGBColor(59, 130, 246),
                'background': RGBColor(253, 242, 248),
                'text': RGBColor(131, 24, 67)
            },
            'corporate_navy': {
                'primary': RGBColor(30, 58, 138),
                'secondary': RGBColor(147, 197, 253),
                'accent': RGBColor(245, 158, 11),
                'background': RGBColor(248, 250, 252),
                'text': RGBColor(30, 41, 59)
            },
            'creative_magenta': {
                'primary': RGBColor(192, 38, 211),
                'secondary': RGBColor(233, 213, 255),
                'accent': RGBColor(16, 185, 129),
                'background': RGBColor(253, 244, 255),
                'text': RGBColor(112, 26, 117)
            },
            'deep_crimson': {
                'primary': RGBColor(220, 38, 127),
                'secondary': RGBColor(252, 165, 165),
                'accent': RGBColor(34, 197, 94),
                'background': RGBColor(254, 242, 242),
                'text': RGBColor(127, 29, 29)
            },
            'royal_indigo': {
                'primary': RGBColor(79, 70, 229),
                'secondary': RGBColor(165, 180, 252),
                'accent': RGBColor(251, 191, 36),
                'background': RGBColor(238, 242, 255),
                'text': RGBColor(30, 27, 75)
            },
            'emerald_mint': {
                'primary': RGBColor(5, 150, 105),
                'secondary': RGBColor(110, 231, 183),
                'accent': RGBColor(249, 115, 22),
                'background': RGBColor(236, 253, 245),
                'text': RGBColor(6, 78, 59)
            },
            'golden_amber': {
                'primary': RGBColor(217, 119, 6),
                'secondary': RGBColor(254, 215, 170),
                'accent': RGBColor(168, 85, 247),
                'background': RGBColor(255, 251, 235),
                'text': RGBColor(120, 53, 15)
            },
            'steel_slate': {
                'primary': RGBColor(71, 85, 105),
                'secondary': RGBColor(203, 213, 225),
                'accent': RGBColor(239, 68, 68),
                'background': RGBColor(248, 250, 252),
                'text': RGBColor(15, 23, 42)
            },
            'cosmic_violet': {
                'primary': RGBColor(124, 58, 237),
                'secondary': RGBColor(196, 181, 253),
                'accent': RGBColor(34, 197, 94),
                'background': RGBColor(245, 243, 255),
                'text': RGBColor(46, 16, 101)
            },
            'cherry_blossom': {
                'primary': RGBColor(244, 63, 94),
                'secondary': RGBColor(252, 231, 243),
                'accent': RGBColor(59, 130, 246),
                'background': RGBColor(255, 241, 242),
                'text': RGBColor(136, 19, 55)
            },
            'arctic_cyan': {
                'primary': RGBColor(6, 182, 212),
                'secondary': RGBColor(165, 243, 252),
                'accent': RGBColor(251, 146, 60),
                'background': RGBColor(236, 254, 255),
                'text': RGBColor(22, 78, 99)
            },
            'sunset_coral': {
                'primary': RGBColor(251, 113, 133),
                'secondary': RGBColor(254, 205, 211),
                'accent': RGBColor(16, 185, 129),
                'background': RGBColor(255, 228, 230),
                'text': RGBColor(159, 18, 57)
            },
            'midnight_blue': {
                'primary': RGBColor(30, 64, 175),
                'secondary': RGBColor(147, 197, 253),
                'accent': RGBColor(245, 158, 11),
                'background': RGBColor(239, 246, 255),
                'text': RGBColor(23, 37, 84)
            },
            'forest_moss': {
                'primary': RGBColor(22, 101, 52),
                'secondary': RGBColor(187, 247, 208),
                'accent': RGBColor(239, 68, 68),
                'background': RGBColor(240, 253, 244),
                'text': RGBColor(14, 59, 30)
            },
            'lavender_dream': {
                'primary': RGBColor(147, 51, 234),
                'secondary': RGBColor(221, 214, 254),
                'accent': RGBColor(251, 146, 60),
                'background': RGBColor(250, 245, 255),
                'text': RGBColor(88, 28, 135)
            },
            'bronze_gold': {
                'primary': RGBColor(180, 83, 9),
                'secondary': RGBColor(253, 186, 116),
                'accent': RGBColor(168, 85, 247),
                'background': RGBColor(255, 247, 237),
                'text': RGBColor(154, 52, 18)
            },
            'ocean_depth': {
                'primary': RGBColor(15, 118, 110),
                'secondary': RGBColor(153, 246, 228),
                'accent': RGBColor(251, 113, 133),
                'background': RGBColor(240, 253, 250),
                'text': RGBColor(19, 78, 74)
            },
            'ruby_wine': {
                'primary': RGBColor(190, 18, 60),
                'secondary': RGBColor(252, 165, 165),
                'accent': RGBColor(34, 197, 94),
                'background': RGBColor(255, 228, 230),
                'text': RGBColor(127, 29, 29)
            }
        }
        
        self.decoration_styles = [
            'geometric_modern', 'organic_flow', 'minimal_lines', 'dynamic_shapes',
            'abstract_art', 'tech_grid', 'nature_inspired'
        ]
        
        self.current_scheme = None
        self.slide_counter = 0
    
    def _get_safe_position(self, slide, min_width, max_width, min_height, max_height, decoration_type="normal"):
        text_areas = self._get_text_areas(slide)
        max_attempts = 100
        
        safe_zones = self._get_safe_zones(text_areas, decoration_type)
        
        for attempt in range(max_attempts):
            width = Inches(random.uniform(min_width, max_width))
            height = Inches(random.uniform(min_height, max_height))
            
            zone = random.choice(safe_zones)
            x = Inches(random.uniform(zone['x_min'], zone['x_max'] - width.inches))
            y = Inches(random.uniform(zone['y_min'], zone['y_max'] - height.inches))
            
            if x.inches >= 0 and y.inches >= 0 and (x.inches + width.inches) <= 10 and (y.inches + height.inches) <= 7.5:
                if not self._overlaps_with_text(x, y, width, height, text_areas):
                    return x, y, width, height
        
        fallback_zone = safe_zones[0]
        fallback_width = Inches(random.uniform(min_width, max_width))
        fallback_height = Inches(random.uniform(min_height, max_height))
        return Inches(max(0, fallback_zone['x_min'])), Inches(max(0, fallback_zone['y_min'])), fallback_width, fallback_height
    
    def _get_text_areas(self, slide):
        text_areas = []
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text_areas.append({
                    'left': shape.left.inches,
                    'top': shape.top.inches,
                    'right': shape.left.inches + shape.width.inches,
                    'bottom': shape.top.inches + shape.height.inches
                })
        return text_areas
    
    def _get_safe_zones(self, text_areas, decoration_type):
        slide_width = 10
        slide_height = 7.5
        
        if decoration_type == "background":
            return [{'x_min': 0.1, 'x_max': 9.9, 'y_min': 0.1, 'y_max': 7.4}]
        
        zones = [
            {'x_min': 7.5, 'x_max': 9.8, 'y_min': 0.2, 'y_max': 2.5},
            {'x_min': 8.0, 'x_max': 9.8, 'y_min': 2.8, 'y_max': 4.5},
            {'x_min': 7.2, 'x_max': 9.8, 'y_min': 5.0, 'y_max': 7.2},
            {'x_min': 0.2, 'x_max': 2.0, 'y_min': 0.2, 'y_max': 1.5},
            {'x_min': 0.2, 'x_max': 1.8, 'y_min': 6.5, 'y_max': 7.3}
        ]
        
        safe_zones = []
        for zone in zones:
            zone_clear = True
            for area in text_areas:
                if not (zone['x_max'] < area['left'] or 
                       zone['x_min'] > area['right'] or 
                       zone['y_max'] < area['top'] or 
                       zone['y_min'] > area['bottom']):
                    zone_clear = False
                    break
            if zone_clear and (zone['x_max'] - zone['x_min']) > 0.5 and (zone['y_max'] - zone['y_min']) > 0.5:
                safe_zones.append(zone)
        
        if not safe_zones:
            safe_zones = [{'x_min': 8.5, 'x_max': 9.8, 'y_min': 6.0, 'y_max': 7.2}]
        
        return safe_zones
    
    def _overlaps_with_text(self, x, y, width, height, text_areas):
        buffer = Inches(0.4)
        decoration_left = x - buffer
        decoration_right = x + width + buffer
        decoration_top = y - buffer
        decoration_bottom = y + height + buffer
        
        for area in text_areas:
            if not (decoration_right < area['left'] or 
                   decoration_left > area['right'] or 
                   decoration_bottom < area['top'] or 
                   decoration_top > area['bottom']):
                return True
        return False
    
    def choose_color_scheme(self, presentation_title: str) -> dict:
        title_lower = presentation_title.lower()
        
        business_keywords = self.loc.t("business_keywords") + ["business", "corporate", "finance", "company"]
        if any(word in title_lower for word in business_keywords):
            return random.choice([self.color_schemes['corporate_navy'], self.color_schemes['steel_slate'], self.color_schemes['midnight_blue']])
        
        creative_keywords = self.loc.t("creative_keywords") + ["creative", "design", "art", "artistic"]
        if any(word in title_lower for word in creative_keywords):
            return random.choice([self.color_schemes['creative_magenta'], self.color_schemes['lavender_dream'], self.color_schemes['cherry_blossom']])
        
        nature_keywords = self.loc.t("nature_keywords") + ["nature", "eco", "green", "forest", "plant"]
        if any(word in title_lower for word in nature_keywords):
            return random.choice([self.color_schemes['forest_green'], self.color_schemes['emerald_mint'], self.color_schemes['forest_moss']])
        
        tech_keywords = self.loc.t("tech_keywords") + ["IT", "tech", "digital", "computer", "internet"]
        if any(word in title_lower for word in tech_keywords):
            return random.choice([self.color_schemes['modern_blue'], self.color_schemes['arctic_cyan'], self.color_schemes['cosmic_violet']])
        
        medical_keywords = self.loc.t("medical_keywords") + ["medicine", "health", "treatment", "doctor"]
        if any(word in title_lower for word in medical_keywords):
            return random.choice([self.color_schemes['ocean_teal'], self.color_schemes['arctic_cyan'], self.color_schemes['emerald_mint']])
        
        education_keywords = self.loc.t("education_keywords") + ["education", "science", "study", "research"]
        if any(word in title_lower for word in education_keywords):
            return random.choice([self.color_schemes['royal_indigo'], self.color_schemes['elegant_purple'], self.color_schemes['cosmic_violet']])
        
        energy_keywords = self.loc.t("energy_keywords") + ["energy", "industry", "production", "factory"]
        if any(word in title_lower for word in energy_keywords):
            return random.choice([self.color_schemes['golden_amber'], self.color_schemes['bronze_gold'], self.color_schemes['warm_orange']])
        
        beauty_keywords = self.loc.t("beauty_keywords") + ["love", "beauty", "fashion", "style"]
        if any(word in title_lower for word in beauty_keywords):
            return random.choice([self.color_schemes['sunset_pink'], self.color_schemes['cherry_blossom'], self.color_schemes['sunset_coral']])
        
        sport_keywords = self.loc.t("sport_keywords") + ["sport", "fitness", "active", "training"]
        if any(word in title_lower for word in sport_keywords):
            return random.choice([self.color_schemes['deep_crimson'], self.color_schemes['ruby_wine'], self.color_schemes['warm_orange']])
        
        travel_keywords = self.loc.t("travel_keywords") + ["ocean", "sea", "travel", "vacation"]
        if any(word in title_lower for word in travel_keywords):
            return random.choice([self.color_schemes['ocean_depth'], self.color_schemes['ocean_teal'], self.color_schemes['arctic_cyan']])
        
        return random.choice(list(self.color_schemes.values()))
    
    def add_slide_decoration(self, slide, slide_index: int, total_slides: int, slide_type: str = "content", presentation_title: str = ""):
        if self.current_scheme is None:
            self.current_scheme = self.choose_color_scheme(presentation_title)
        
        self.slide_counter += 1
        decoration_style = self._choose_decoration_style(slide_index, total_slides, slide_type)
        
        if decoration_style == "geometric_modern":
            self._add_geometric_modern(slide)
        elif decoration_style == "organic_flow":
            self._add_organic_flow(slide)
        elif decoration_style == "minimal_lines":
            self._add_minimal_lines(slide)
        elif decoration_style == "dynamic_shapes":
            self._add_dynamic_shapes(slide)
        elif decoration_style == "abstract_art":
            self._add_abstract_art(slide)
        elif decoration_style == "tech_grid":
            self._add_tech_grid(slide)
        elif decoration_style == "nature_inspired":
            self._add_nature_inspired(slide)
    
    def _choose_decoration_style(self, slide_index: int, total_slides: int, slide_type: str) -> str:
        if slide_index == 0:
            return random.choice(['geometric_modern', 'abstract_art', 'minimal_lines'])
        elif slide_type == "section":
            return random.choice(['minimal_lines', 'dynamic_shapes', 'organic_flow'])
        else:
            used_styles = getattr(self, '_used_styles', [])
            available_styles = [s for s in self.decoration_styles if s not in used_styles[-3:]]
            
            if not available_styles:
                available_styles = self.decoration_styles
                self._used_styles = []
            
            chosen_style = random.choice(available_styles)
            
            if not hasattr(self, '_used_styles'):
                self._used_styles = []
            self._used_styles.append(chosen_style)
            
            return chosen_style
    
    def _add_geometric_modern(self, slide):
        shapes_count = random.randint(2, 4)
        for i in range(shapes_count):
            shape_type = random.choice([MSO_SHAPE.RECTANGLE, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.HEXAGON])
            
            x, y, width, height = self._get_safe_position(slide, 0.6, 1.5, 0.6, 1.5, "small")
            
            shape = slide.shapes.add_shape(shape_type, x, y, width, height)
            
            if i % 2 == 0:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = self.current_scheme['primary']
                shape.line.fill.background()
            else:
                shape.fill.background()
                line = shape.line
                line.color.rgb = self.current_scheme['secondary']
                line.width = Pt(2)
            
            shape.rotation = random.randint(-20, 20)
    
    def _add_organic_flow(self, slide):
        for i in range(random.randint(2, 3)):
            x, y, width, height = self._get_safe_position(slide, 0.8, 2.0, 0.4, 1.5, "medium")
            
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, width, height)
            
            fill = shape.fill
            fill.solid()
            base_color = self.current_scheme['secondary']
            fill.fore_color.rgb = RGBColor(
                min(255, max(50, base_color[0] + random.randint(-30, 30))),
                min(255, max(50, base_color[1] + random.randint(-30, 30))),
                min(255, max(50, base_color[2] + random.randint(-30, 30)))
            )
            
            shape.line.fill.background()
            shape.rotation = random.randint(-30, 30)
    
    def _add_minimal_lines(self, slide):
        line_count = random.randint(3, 5)
        for i in range(line_count):
            safe_zones = self._get_safe_zones(self._get_text_areas(slide), "line")
            if safe_zones:
                zone = random.choice(safe_zones)
                
            if random.choice([True, False]):
                line = slide.shapes.add_connector(
                    1,
                        Inches(zone['x_min']),
                        Inches(random.uniform(zone['y_min'], zone['y_max'])),
                        Inches(zone['x_max']),
                        Inches(random.uniform(zone['y_min'], zone['y_max']))
                )
            else:
                line = slide.shapes.add_connector(
                    1,
                        Inches(random.uniform(zone['x_min'], zone['x_max'])),
                        Inches(zone['y_min']),
                        Inches(random.uniform(zone['x_min'], zone['x_max'])),
                        Inches(zone['y_max'])
                )
            
                line.line.color.rgb = self.current_scheme['primary']
                line.line.width = Pt(random.randint(1, 3))
    
    def _add_dynamic_shapes(self, slide):
        shapes = [MSO_SHAPE.DIAMOND, MSO_SHAPE.PENTAGON, MSO_SHAPE.HEXAGON, MSO_SHAPE.ROUNDED_RECTANGLE]
        
        for i in range(random.randint(2, 4)):
            x, y, width, height = self._get_safe_position(slide, 0.5, 1.2, 0.5, 1.2, "small")
            
            shape = slide.shapes.add_shape(random.choice(shapes), x, y, width, height)
            
            if i % 3 == 0:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = self.current_scheme['accent']
                shape.line.fill.background()
            else:
                shape.fill.background()
                line = shape.line
                line.color.rgb = self.current_scheme['primary']
                line.width = Pt(1)
            
            shape.rotation = random.randint(0, 180)
    
    def _add_gradient_waves(self, slide):
        text_areas = self._get_text_areas(slide)
        
        safe_y_position = 6.0
        for area in text_areas:
            if area['bottom'] > safe_y_position:
                safe_y_position = area['bottom'] + 0.3
        
        if safe_y_position > 6.5:
            safe_y_position = 0.2
            wave_height = 1.0
        else:
            wave_height = min(1.5, 7.5 - safe_y_position - 0.2)
        
        wave_shape = slide.shapes.add_shape(
            MSO_SHAPE.WAVE,
            Inches(0), Inches(safe_y_position),
            Inches(10), Inches(wave_height)
        )
        
        fill = wave_shape.fill
        fill.solid()
        fill.fore_color.rgb = self.current_scheme['background']
        
        line = wave_shape.line
        line.color.rgb = self.current_scheme['secondary']
        line.width = Pt(2)
    
    def _add_abstract_art(self, slide):
        for i in range(random.randint(3, 5)):
            shape_type = random.choice([MSO_SHAPE.OVAL, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.DIAMOND])
            
            x, y, width, height = self._get_safe_position(slide, 0.4, 1.0, 0.4, 1.0, "small")
            
            shape = slide.shapes.add_shape(shape_type, x, y, width, height)
            
            colors = [self.current_scheme['primary'], self.current_scheme['secondary'], self.current_scheme['accent']]
            
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = random.choice(colors)
            
            shape.line.fill.background()
            shape.rotation = random.randint(0, 180)
    
    def _add_tech_grid(self, slide):
        grid_size = 0.3
        for x in range(int(10 / grid_size)):
            for y in range(int(7.5 / grid_size)):
                if random.random() < 0.1:
                    dot = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(x * grid_size),
                        Inches(y * grid_size),
                        Inches(0.05),
                        Inches(0.05)
                    )
                    
                    fill = dot.fill
                    fill.solid()
                    fill.fore_color.rgb = self.current_scheme['secondary']
                    dot.line.fill.background()
    
    def _add_nature_inspired(self, slide):
        leaf_shapes = [MSO_SHAPE.OVAL, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.TEAR]
        
        for i in range(random.randint(2, 4)):
            x, y, width, height = self._get_safe_position(slide, 0.5, 1.0, 0.6, 1.5, "medium")
            
            shape = slide.shapes.add_shape(random.choice(leaf_shapes), x, y, width, height)
            
            green_variations = [
                RGBColor(34, 197, 94),
                RGBColor(22, 163, 74),
                RGBColor(21, 128, 61),
                RGBColor(134, 239, 172)
            ]
            
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = random.choice(green_variations)
            
            shape.line.fill.background()
            shape.rotation = random.randint(-20, 20)
    
    def apply_advanced_text_formatting(self, text_frame, slide_type: str = "content", content: str = ""):
        if self.current_scheme is None:
            self.current_scheme = self.choose_color_scheme("")
            
        if slide_type == "title":
            self._format_title_advanced(text_frame)
        elif slide_type == "section":
            self._format_section_advanced(text_frame)
        elif slide_type == "content":
            self._format_content_advanced(text_frame, content)
        elif slide_type == "list":
            self._format_list_advanced(text_frame, content)
    
    def _format_title_advanced(self, text_frame):
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = 'Montserrat'
                font.size = Pt(48)
                font.bold = True
                font.color.rgb = self.current_scheme['primary']
                
                if any(word in run.text.lower() for word in ['важн', 'ключев', 'главн', 'important', 'key', 'main']):
                    font.color.rgb = self.current_scheme['accent']
    
    def _format_section_advanced(self, text_frame):
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = 'Montserrat'
                font.size = Pt(36)
                font.bold = True
                font.color.rgb = self.current_scheme['primary']
    
    def _format_content_advanced(self, text_frame, content: str):
        sentences = content.split('.')
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(12)
            
            if ':' in paragraph.text:
                parts = paragraph.text.split(':', 1)
                if len(parts) >= 2:
                    paragraph.clear()
                    
                    title_run = paragraph.add_run()
                    title_run.text = parts[0] + ':'
                    title_run.font.name = 'Montserrat'
                    title_run.font.size = Pt(18)
                    title_run.font.bold = True
                    title_run.font.color.rgb = self.current_scheme['primary']
                    
                    content_run = paragraph.add_run()
                    content_run.text = parts[1]
                    content_run.font.name = 'Open Sans'
                    content_run.font.size = Pt(16)
                    content_run.font.color.rgb = self.current_scheme['text']
            else:
                for run in paragraph.runs:
                    font = run.font
                    font.name = 'Open Sans'
                    font.size = Pt(16)
                    font.color.rgb = self.current_scheme['text']
                    
                    if any(word in run.text.lower() for word in ['важн', 'ключев', 'главн', 'important', 'key', 'main']):
                        font.bold = True
                        font.color.rgb = self.current_scheme['accent']
    
    def _format_list_advanced(self, text_frame, content: str):
        for paragraph in text_frame.paragraphs:
            paragraph.space_after = Pt(8)
            
            if paragraph.text.strip().startswith(('•', '-', '1.', '2.', '3.')):
                for run in paragraph.runs:
                    if run.text.strip().startswith(('•', '-')):
                        run.font.color.rgb = self.current_scheme['accent']
                        run.font.size = Pt(20)
                        run.font.bold = True
                    else:
                        run.font.name = 'Open Sans'
                        run.font.size = Pt(16)
                        run.font.color.rgb = self.current_scheme['text']
    
    def add_slide_background(self, slide, slide_type: str = "content"):
        if self.current_scheme is None:
            self.current_scheme = self.choose_color_scheme("")
            
        if slide_type == "title":
            self._add_title_background(slide)
        elif random.random() < 0.3:
            self._add_subtle_background(slide)
    
    def _add_title_background(self, slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        
        fill = bg_shape.fill
        fill.solid()
        fill.fore_color.rgb = self.current_scheme['background']
        
        bg_shape.line.fill.background()
        
        bg_shape.element.getparent().remove(bg_shape.element)
        slide.shapes._spTree.insert(2, bg_shape.element)
    
    def _add_subtle_background(self, slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.3), Inches(1.5),
            Inches(9.4), Inches(5.5)
        )
        
        fill = bg_shape.fill
        fill.solid()
        bg_color = self.current_scheme['background']
        fill.fore_color.rgb = RGBColor(250, 250, 252)
        
        line = bg_shape.line
        line.color.rgb = self.current_scheme['secondary']
        line.width = Pt(1)
        
        bg_shape.element.getparent().remove(bg_shape.element)
        slide.shapes._spTree.insert(2, bg_shape.element)
