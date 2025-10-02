import os
import sys
import re
from pathlib import Path
from typing import Optional, List
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from ..models.presentation import Presentation
from ..localization.manager import get_localization_manager
from .decorations import SlideDecorator


def get_resource_path(relative_path):
    if getattr(sys, 'frozen', False):
        return Path(sys._MEIPASS) / relative_path
    else:
        return Path(__file__).parent.parent.parent / relative_path


class PPTXGenerator:
    def __init__(self, output_dir: str = "output"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.decorator = SlideDecorator()
        self.loc = get_localization_manager()
        
    def _create_title_slide(self, pptx: PPTXPresentation, presentation: Presentation, 
                           slide_index: int = 0, total_slides: int = 1) -> None:
        title_slide_layout = pptx.slide_layouts[0]
        slide = pptx.slides.add_slide(title_slide_layout)
        
        self.decorator.add_slide_background(slide, "title")
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        if presentation.title_slide_header:
            title.text = presentation.title_slide_header
        else:
            title.text = presentation.title
            
        if presentation.summary:
            subtitle.text = presentation.summary
        else:
            subtitle.text = self.loc.t("created_with_ai")
            
        self._apply_advanced_title_styling(title.text_frame, subtitle.text_frame)
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "title", presentation.title)
        self.decorator.apply_advanced_text_formatting(title.text_frame, "title")
    
    def _create_section_title_slide(self, pptx: PPTXPresentation, section_title: str, 
                                   slide_index: int, total_slides: int, presentation_title: str = "") -> None:
        section_header_layout = pptx.slide_layouts[2]
        slide = pptx.slides.add_slide(section_header_layout)
        
        title = slide.shapes.title
        title.text = section_title
        
        self._apply_section_styling(title.text_frame)
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "section", presentation_title)
        self.decorator.apply_advanced_text_formatting(title.text_frame, "section")
    
    def _create_content_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                             slide_index: int, total_slides: int, presentation_title: str = "") -> None:
        if not self._validate_slide_content(slide_content):
            slide_content = f"{self.loc.t('slide_content_default')} '{slide_title}'"
        
        if self._is_table_content(slide_content):
            table_data = self._parse_table_data(slide_content)
            if table_data and len(table_data) >= 2:
                slide_layout = pptx.slide_layouts[5]
                pptx_slide = pptx.slides.add_slide(slide_layout)
                self._create_table_slide(pptx_slide, slide_title, table_data)
                return
            else:
                slide_content = f"{self.loc.t('insufficient_table_data')}. Содержимое: {slide_title}"
        
        slide_type = self._determine_slide_type(slide_content)
        
        if slide_type == "list":
            self._create_enhanced_list_slide(pptx, slide_title, slide_content, slide_index, total_slides, presentation_title)
        elif slide_type == "comparison":
            self._create_comparison_slide(pptx, slide_title, slide_content, slide_index, total_slides, presentation_title)
        elif slide_type == "highlight":
            self._create_highlight_slide(pptx, slide_title, slide_content, slide_index, total_slides, presentation_title)
        elif slide_type == "process":
            self._create_process_slide(pptx, slide_title, slide_content, slide_index, total_slides, presentation_title)
        else:
            self._create_modern_content_slide(pptx, slide_title, slide_content, slide_index, total_slides, presentation_title)
    
    def _determine_slide_type(self, content: str) -> str:
        return "content"
    
    def _create_content_shape(self, slide):
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(0.5),
            Inches(9.4), Inches(6.5)
        )
        
        fill = content_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(248, 250, 252)
        
        line = content_shape.line
        line.color.rgb = self.decorator.current_scheme['secondary']
        line.width = Pt(2)
        
        content_frame = content_shape.text_frame
        content_frame.margin_left = Inches(0.3)
        content_frame.margin_right = Inches(0.3)
        content_frame.margin_top = Inches(0.2)
        content_frame.margin_bottom = Inches(0.4)
        content_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        content_frame.word_wrap = True
        
        return content_frame
    
    def _create_enhanced_list_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                                   slide_index: int, total_slides: int, presentation_title: str) -> None:
        blank_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_layout)
        
        self.decorator.add_slide_background(slide, "content")
        content_frame = self._create_content_shape(slide)
        
        full_content = f"{slide_title}\n\n{slide_content}"
        truncated_content = self._truncate_content_for_slide(full_content, "list")
        self._add_formatted_content_to_frame(content_frame, truncated_content)
        
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "list", presentation_title)
    
    def _create_comparison_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                                slide_index: int, total_slides: int, presentation_title: str) -> None:
        blank_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_layout)
        
        self.decorator.add_slide_background(slide, "content")
        content_frame = self._create_content_shape(slide)
        
        full_content = f"{slide_title}\n\n{slide_content}"
        truncated_content = self._truncate_content_for_slide(full_content, "comparison")
        self._add_formatted_content_to_frame(content_frame, truncated_content)
        
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "comparison", presentation_title)
    
    def _create_highlight_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                               slide_index: int, total_slides: int, presentation_title: str) -> None:
        blank_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_layout)
        
        content_frame = self._create_content_shape(slide)
        
        full_content = f"{slide_title}\n\n{slide_content}"
        truncated_content = self._truncate_content_for_slide(full_content, "highlight")
        self._add_formatted_content_to_frame(content_frame, truncated_content)
        
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "content", presentation_title)
    
    def _create_process_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                             slide_index: int, total_slides: int, presentation_title: str) -> None:
        blank_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_layout)
        
        self.decorator.add_slide_background(slide, "content")
        content_frame = self._create_content_shape(slide)
        
        full_content = f"{slide_title}\n\n{slide_content}"
        truncated_content = self._truncate_content_for_slide(full_content, "process")
        self._add_formatted_content_to_frame(content_frame, truncated_content)
        
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "process", presentation_title)
    
    def _create_modern_content_slide(self, pptx: PPTXPresentation, slide_title: str, slide_content: str, 
                                    slide_index: int, total_slides: int, presentation_title: str) -> None:
        blank_layout = pptx.slide_layouts[6]
        slide = pptx.slides.add_slide(blank_layout)
        
        self.decorator.add_slide_background(slide, "content")
        content_frame = self._create_content_shape(slide)
        
        full_content = f"{slide_title}\n\n{slide_content}"
        truncated_content = self._truncate_content_for_slide(full_content, "content")
        self._add_formatted_content_to_frame(content_frame, truncated_content)
        
        self.decorator.add_slide_decoration(slide, slide_index, total_slides, "content", presentation_title)
    
    def _extract_enhanced_list_items(self, content: str) -> List[str]:
        lines = content.split('\n')
        items = []
        current_item = ""
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith(('•', '-', '1.', '2.', '3.', '4.', '5.')):
                if current_item:
                    items.append(current_item)
                current_item = re.sub(r'^[•\-\d\.]\s*', '', line).strip()
            else:
                if current_item:
                    current_item += f" {line}"
                else:
                    current_item = line
        
        if current_item:
            items.append(current_item)
        
        if not items:
            sentences = content.split('.')
            items = [s.strip() for s in sentences if s.strip()][:6]
        
        return items[:6]
    
    def _split_comparison_content(self, content: str) -> dict:
        content_lower = content.lower()
        
        if 'преимущества' in content_lower and 'недостатки' in content_lower:
            parts = content.split('недостатки', 1)
            left = parts[0].replace('преимущества', '').strip()
            right = parts[1].strip() if len(parts) > 1 else ""
        elif 'плюсы' in content_lower and 'минусы' in content_lower:
            parts = content.split('минусы', 1)
            left = parts[0].replace('плюсы', '').strip()
            right = parts[1].strip() if len(parts) > 1 else ""
        else:
            mid_point = len(content) // 2
            left = content[:mid_point].strip()
            right = content[mid_point:].strip()
        
        return {'left': left, 'right': right}
    
    def _extract_process_steps(self, content: str) -> List[str]:
        steps = []
        sentences = content.split('.')
        
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(sentence) > 10:
                steps.append(sentence)
        
        return steps[:5]
    
    def _apply_advanced_title_styling(self, title_frame, subtitle_frame) -> None:
        title_frame.word_wrap = True
        title_frame.auto_size = None
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        font = run.font
        font.name = 'Montserrat'
        font.size = Pt(52)
        font.bold = True
        font.color.rgb = self.decorator.current_scheme['primary']
        
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_p.runs[0] if subtitle_p.runs else subtitle_p.add_run()
        subtitle_font = subtitle_run.font
        subtitle_font.name = 'Open Sans'
        subtitle_font.size = Pt(22)
        subtitle_font.color.rgb = self.decorator.current_scheme['text']
    
    def _apply_section_styling(self, text_frame) -> None:
        text_frame.word_wrap = True
        text_frame.auto_size = None
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        font = run.font
        font.name = 'Montserrat'
        font.size = Pt(40)
        font.bold = True
        font.color.rgb = self.decorator.current_scheme['primary']
    
    def _apply_slide_title_styling(self, text_frame) -> None:
        text_frame.word_wrap = True
        text_frame.auto_size = None
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        font = run.font
        font.name = 'Montserrat'
        font.size = Pt(36)
        font.bold = True
        font.color.rgb = self.decorator.current_scheme['primary']
    
    def _apply_enhanced_list_styling(self, paragraph, index: int) -> None:
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(14)
        paragraph.line_spacing = 1.2
        
        text = paragraph.text
        paragraph.clear()
        
        if text.startswith('•'):
            bullet_run = paragraph.add_run()
            bullet_run.text = '• '
            bullet_run.font.name = 'Arial'
            bullet_run.font.size = Pt(18)
            bullet_run.font.color.rgb = self.decorator.current_scheme['accent']
            bullet_run.font.bold = True
            
            content_run = paragraph.add_run()
            content_run.text = text[2:].strip()
            content_run.font.name = 'Open Sans'
            content_run.font.size = Pt(16)
            content_run.font.color.rgb = self.decorator.current_scheme['text']
        else:
            content_run = paragraph.add_run()
            content_run.text = text
            content_run.font.name = 'Open Sans'
            content_run.font.size = Pt(16)
            content_run.font.color.rgb = self.decorator.current_scheme['text']
    
    def _apply_comparison_styling(self, text_frame, side: str) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(12)
            paragraph.line_spacing = 1.2
            for run in paragraph.runs:
                font = run.font
                font.name = 'Open Sans'
                font.size = Pt(16)
                if side == "positive":
                    font.color.rgb = self.decorator.current_scheme['primary']
                else:
                    font.color.rgb = self.decorator.current_scheme['text']
    
    def _apply_highlight_styling(self, text_frame) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = 'Montserrat'
                font.size = Pt(20)
                font.bold = True
                font.color.rgb = RGBColor(255, 255, 255)
    
    def _apply_process_step_styling(self, text_frame, is_primary: bool) -> None:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = 'Open Sans'
                font.size = Pt(14)
                font.bold = True
                if is_primary:
                    font.color.rgb = RGBColor(255, 255, 255)
                else:
                    font.color.rgb = self.decorator.current_scheme['text']
    
    def _is_table_content(self, content: str) -> bool:
        return content.strip().startswith('TABLE|')
    
    def _parse_table_data(self, content: str) -> list:
        if not self._is_table_content(content):
            return []
        
        lines = content.strip().split('\n')
        table_data = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('TABLE|'):
                headers = line[6:].split('|')
                table_data.append([header.strip() for header in headers if header.strip()])
            else:
                row = line.split('|')
                if len(row) >= 2:
                    table_data.append([cell.strip() for cell in row if cell.strip()])
        
        if len(table_data) < 2 or len(table_data[0]) < 2:
            return []
            
        return table_data
    
    def _create_table_slide(self, slide, title: str, table_data: list) -> None:
        try:
            if not table_data or len(table_data) < 2:
                self._add_error_text(slide, self.loc.t("insufficient_table_data"))
                return
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.text = title
            self._apply_slide_title_styling(title_frame)
        
            rows = len(table_data)
            cols = len(table_data[0]) if table_data[0] else 1
            
            if rows > 6:
                rows = 6
                table_data = table_data[:6]
            if cols > 5:
                cols = 5
                table_data = [row[:5] if len(row) >= 5 else row + [''] * (5 - len(row)) for row in table_data]
            
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(5)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    if j < cols and i < rows:
                        cell = table.cell(i, j)
                        cell.text = str(cell_data) if cell_data else ""
                        
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.name = 'Open Sans'
                                
                                if i == 0:
                                    run.font.size = Pt(14)
                                    run.font.bold = True
                                    run.font.color.rgb = self.decorator.current_scheme['primary']
                                else:
                                    run.font.size = Pt(12)
                                    run.font.color.rgb = self.decorator.current_scheme['text']
                        
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = self.decorator.current_scheme['secondary']
            
            self.decorator.add_slide_background(slide)
            self._add_slide_number(slide)
            
        except Exception as e:
            print(f"{self.loc.t('table_creation_error')}: {e}")
            self._add_error_text(slide, self.loc.t("table_creation_failed"))
    
    def _add_formatted_content_to_frame(self, content_frame, content: str) -> None:
        if self._is_table_content(content):
            return
        
        content_frame.clear()
        lines = content.split('\n')
        is_first_line = True
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            p = content_frame.add_paragraph()
            
            if is_first_line:
                run = p.add_run()
                run.text = line
                run.font.name = 'Open Sans'
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = self.decorator.current_scheme['primary']
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)
                is_first_line = False
            elif line.startswith('•'):
                p.level = 0
                bullet_run = p.add_run()
                bullet_run.text = '• '
                bullet_run.font.name = 'Open Sans'
                bullet_run.font.size = Pt(16)
                bullet_run.font.color.rgb = self.decorator.current_scheme['primary']
                bullet_run.font.bold = True
                
                text_run = p.add_run()
                text_run.text = line[1:].strip()
                text_run.font.name = 'Open Sans'
                text_run.font.size = Pt(16)
                text_run.font.color.rgb = self.decorator.current_scheme['text']
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(8)
            else:
                run = p.add_run()
                run.text = line
                run.font.name = 'Open Sans'
                run.font.size = Pt(16)
                run.font.color.rgb = self.decorator.current_scheme['text']
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(8)
            
            p.line_spacing = 1.2

    def _validate_slide_content(self, content: str) -> bool:
        if not content or len(content.strip()) < 10:
            return False
        
        content_lower = content.lower()
        invalid_indicators = [
            self.loc.t("content_for_slide"),
            "content for slide",
            self.loc.t("placeholder"),
            self.loc.t("placeholder"),
            self.loc.t("enter_text"),
            self.loc.t("add_content"),
            self.loc.t("text_will_be_here"),
            self.loc.t("text_for_slide")
        ]
        
        return not any(indicator in content_lower for indicator in invalid_indicators)
    
    async def _regenerate_slide_content(self, slide_title: str, section_title: str, language: str = None) -> str:
        if language is None:
            language = self.loc.t('language_russian')
        from ..services.ai_service import AIService
        from ..localization.manager import get_localization_manager
        
        loc = get_localization_manager()
        api_key = self._get_api_key()
        
        if not api_key:
            return f"{loc.t('slide_content_default')} '{slide_title}'"
        
        ai_service = AIService(api_key)
        return await ai_service.generate_slide_content(slide_title, section_title, language)
    
    def _get_api_key(self) -> str:
        try:
            with open("config/api_key.txt", "r", encoding="utf-8") as f:
                return f.read().strip()
        except:
            return ""
    
    def _truncate_content_for_slide(self, content: str, slide_type: str) -> str:
        max_chars = {
            "content": 400,
            "list": 300,
            "comparison": 350,
            "highlight": 250,
            "process": 400
        }
        
        max_length = max_chars.get(slide_type, 1200)
        
        if len(content) <= max_length:
            return content
        
        lines = content.split('\n')
        truncated_lines = []
        current_length = 0
        
        for line in lines:
            if current_length + len(line) <= max_length:
                truncated_lines.append(line)
                current_length += len(line) + 1
            else:
                break
        
        if truncated_lines:
            return '\n'.join(truncated_lines)
        
        truncated = content[:max_length]
        last_period = truncated.rfind('.')
        last_exclamation = truncated.rfind('!')
        last_question = truncated.rfind('?')
        
        last_sentence_end = max(last_period, last_exclamation, last_question)
        
        if last_sentence_end > max_length * 0.7:
            return truncated[:last_sentence_end + 1]
        
        return truncated + "..."
    
    def _sanitize_filename(self, filename: str) -> str:
        import re
        filename = filename[:50]
        invalid_chars = '<>:"/\\|?*«»„"'
        for char in invalid_chars:
            filename = filename.replace(char, '_')
        filename = re.sub(r'[^\w\s-]', '_', filename)
        filename = re.sub(r'\s+', '_', filename)
        filename = re.sub(r'_+', '_', filename)
        return filename.strip('_') or "Presentation"
    
    def _generate_smart_filename(self, presentation: Presentation) -> Optional[str]:
        try:
            api_key = self._get_api_key()
            if not api_key:
                return None
            
            ai_service = AIService(api_key)
            import asyncio
            filename = asyncio.run(ai_service.generate_filename(
                presentation.title, 
                presentation.language
            ))
            
            if filename and filename != "Presentation":
                return f"{filename}.pptx"
            return None
        except Exception:
            return None
    
    def generate_pptx(self, presentation: Presentation, filename: Optional[str] = None) -> str:
        if not filename:
            filename = self._generate_smart_filename(presentation)
            if not filename:
                safe_title = self._sanitize_filename(presentation.title)
                if len(safe_title) > 30:
                    words = safe_title.split('_')[:3]
                    safe_title = '_'.join(words)
                filename = f"{safe_title}.pptx"
        
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        if len(filename) > 100:
            filename = "Presentation.pptx"
            
        output_path = self.output_dir / filename
        pptx = PPTXPresentation()
        
        total_slides = 1
        if presentation.sections:
            total_slides += len(presentation.sections)
            for section in presentation.sections:
                if section.slides:
                    total_slides += len(section.slides)
        
        slide_index = 0
        
        self._create_title_slide(pptx, presentation, slide_index, total_slides)
        slide_index += 1
        
        if presentation.sections:
            for section in presentation.sections:
                if section.slides and len(section.slides) > 0:
                    self._create_section_title_slide(pptx, section.title, slide_index, total_slides, presentation.title)
                    slide_index += 1
                    
                    for slide in section.slides:
                        self._create_content_slide(pptx, slide.title, slide.content, slide_index, total_slides, presentation.title)
                        slide_index += 1
        
        pptx.save(str(output_path))
        return str(output_path.absolute())
    
    def list_presentations(self) -> list:
        return [f.name for f in self.output_dir.glob("*.pptx")]
    
    def delete_presentation(self, filename: str) -> bool:
        try:
            file_path = self.output_dir / filename
            if file_path.exists() and file_path.suffix == '.pptx':
                file_path.unlink()
                return True
        except Exception:
            pass
        return False
    
    def _set_multiline_text(self, text_frame, text: str, is_bullet: bool = False, bullet_char: str = '•'):
        text_frame.clear()
        
        lines = [ln for ln in text.split('\n') if ln.strip()]
        if not lines:
            return
        
        for i, ln in enumerate(lines):
            p = text_frame.add_paragraph()
            
            if ln.strip().startswith('•'):
                p.level = 1
                p.space_before = Pt(4)
                p.space_after = Pt(4)
            elif ln.strip().endswith(':'):
                p.level = 0
                p.space_before = Pt(6)
                p.space_after = Pt(3)
            else:
                p.level = 0
                p.space_before = Pt(3)
                p.space_after = Pt(6)
            
            run = p.add_run()
            run.text = ln
            run.font.name = 'Open Sans'
            run.font.size = Pt(15 if ln.strip().startswith('•') else 16)
            run.font.color.rgb = self.decorator.current_scheme['text']
            
            if ln.strip().endswith(':'):
                run.font.bold = True
            
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.2
        return False
    
    def _add_error_text(self, slide, error_message: str) -> None:
        error_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        error_frame = error_box.text_frame
        error_frame.text = error_message
        error_frame.word_wrap = True
        
        for paragraph in error_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Open Sans'
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(200, 50, 50)
                run.font.bold = True
    
    def _add_slide_number(self, slide) -> None:
        slide_number_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.8), Inches(1.5), Inches(0.4))
        slide_number_frame = slide_number_box.text_frame
        slide_number_frame.text = f"{slide.slide_id if hasattr(slide, 'slide_id') else '1'}"
        
        for paragraph in slide_number_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Open Sans'
                run.font.size = Pt(12)
                run.font.color.rgb = self.decorator.current_scheme['text']
